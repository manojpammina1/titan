#!/usr/bin/env python3
"""Build Part 2 PowerPoint — light ocean theme, screenshots on every slide.

Run from repo root:
    docs/part2-deck/.venv/bin/python docs/part2-deck/build-pptx.py

Output:
    docs/part2-deck/Part2-Agentic-Development.pptx
"""
from __future__ import annotations

import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DECK = Path(__file__).resolve().parent
ROOT = DECK.parents[1]
SHOTS = DECK / "screenshots"
OUT = DECK / "Part2-Agentic-Development.pptx"

# ── palette — light deck, ocean / commerce accent (not black) ───────────────

C = {
    "ocean_deep": RGBColor(0x00, 0x5F, 0x8C),
    "ocean_mid": RGBColor(0x00, 0x96, 0xC7),
    "ocean_light": RGBColor(0x90, 0xE0, 0xEF),
    "sky": RGBColor(0xE8, 0xF4, 0xFC),
    "paper": RGBColor(0xFF, 0xFF, 0xFF),
    "mist": RGBColor(0xF1, 0xF7, 0xFB),
    "title": RGBColor(0x0A, 0x3D, 0x5C),
    "body": RGBColor(0x33, 0x41, 0x55),
    "muted": RGBColor(0x64, 0x74, 0x8B),
    "coral": RGBColor(0xFF, 0x7A, 0x45),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "card_border": RGBColor(0xCB, 0xE8, 0xF5),
}


def _load_slides() -> list[dict]:
    spec = importlib.util.spec_from_file_location("build_deck", DECK / "build-deck.py")
    mod = importlib.util.module_from_spec(spec)
    assert spec.loader
    spec.loader.exec_module(mod)
    return mod.SLIDES


def _gradient(shape, c0: RGBColor, c1: RGBColor, angle: float = 90.0) -> None:
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = c0
    stops[0].position = 0.0
    stops[1].color.rgb = c1
    stops[1].position = 1.0


def _solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _set_run(run, *, size: int, bold: bool = False, color: RGBColor = C["body"]) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_bg(slide) -> None:
    """Soft sky gradient + decorative wave bands."""
    w, h = Inches(13.333), Inches(7.5)
    base = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, w, h)
    _gradient(base, C["paper"], C["sky"], 110)
    _no_line(base)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(5.8), w, Inches(1.7))
    _gradient(band, C["ocean_light"], C["sky"], 0)
    band.fill.transparency = 0.35
    _no_line(band)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, w, Inches(0.12))
    _gradient(accent, C["coral"], C["ocean_mid"], 0)
    _no_line(accent)


def _add_header(slide, title: str, slide_num: int, total: int) -> None:
    w = Inches(13.333)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.12), w, Inches(1.05))
    _gradient(header, C["ocean_deep"], C["ocean_mid"], 0)
    _no_line(header)

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(10.5), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=28, bold=True, color=C["white"])

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.85), Inches(0.35), Inches(1.05), Inches(0.45)
    )
    _solid(badge, C["white"])
    badge.fill.transparency = 0.15
    _no_line(badge)
    btf = badge.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"{slide_num}/{total}"
    _set_run(br, size=11, bold=True, color=C["white"])


def _add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(4.35), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        dot = p.add_run()
        dot.text = "●  "
        _set_run(dot, size=13, bold=True, color=C["ocean_mid"])
        body = p.add_run()
        body.text = text
        _set_run(body, size=15, color=C["body"])


def _add_image_card(slide, img_path: Path, left, top, width, height) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    _solid(card, C["paper"])
    card.line.color.rgb = C["card_border"]
    card.line.width = Pt(1.25)
    pad = Inches(0.08)
    slide.shapes.add_picture(
        str(img_path),
        left + pad,
        top + pad,
        width=width - pad * 2,
        height=height - pad * 2,
    )


def _add_footer(slide) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(8), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "RCG Principal Engineer · Part 2 · Agentic Development Strategy"
    _set_run(run, size=10, color=C["muted"])


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    w = Inches(13.333)

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(2.0), w, Inches(3.2))
    _gradient(hero, C["ocean_deep"], C["ocean_mid"], 15)
    hero.fill.transparency = 0.08
    _no_line(hero)

    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.4))
    tp = title.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "Enabling Agentic Development"
    _set_run(tr, size=40, bold=True, color=C["title"])

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.0))
    sp = sub.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "30 / 60 / 90 roadmap · SOX-regulated e-commerce engineering"
    _set_run(sr, size=20, color=C["ocean_deep"])

    tag = slide.shapes.add_textbox(Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.6))
    tgr = tag.text_frame.paragraphs[0].add_run()
    tgr.text = "Principal Engineer case study · Royal Caribbean Group"
    _set_run(tgr, size=14, color=C["muted"])


def build() -> Path:
    slides_data = _load_slides()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)
    total = len(slides_data)

    for s in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide)
        _add_header(slide, s["title"], s["n"], total)
        _add_bullets(slide, s["bullets"])
        _add_footer(slide)

        imgs = [SHOTS / s["shot"]]
        if s.get("shot2"):
            imgs.append(SHOTS / s["shot2"])

        right = Inches(5.15)
        if len(imgs) == 1:
            _add_image_card(slide, imgs[0], right, Inches(1.45), Inches(7.85), Inches(5.35))
        else:
            h_each = Inches(2.58)
            gap = Inches(0.18)
            _add_image_card(slide, imgs[0], right, Inches(1.45), Inches(7.85), h_each)
            _add_image_card(slide, imgs[1], right, Inches(1.45) + h_each + gap, Inches(7.85), h_each)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path} ({path.stat().st_size // 1024} KB)")
